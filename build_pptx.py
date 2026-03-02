#!/usr/bin/env python3
"""
Build CSCE 714 Midterm Presentation.pptx
Updated narrative structure: intro → edge relevance → datasets → prior research → references
"""

import zipfile, io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

ODP_PATH  = 'CSCE 714 Midterm Presentation.odp'
TMPL_PATH = '/Users/kuba/.claude/skills/powerpoint-maker/assets/usc_engineering_powerpoint33.potx'
OUT_PATH  = 'CSCE 714 Midterm Presentation.pptx'

# ── template loader (patch .potx → .pptx content type) ──────────────────────

def _delete_slide(prs, index):
    """Remove a slide by index (python-pptx has no built-in for this)."""
    from pptx.oxml.ns import qn
    xml_slides = prs.slides._sldIdLst
    slide      = prs.slides[index]
    rId        = xml_slides[index].get(qn('r:id'))
    prs.part.drop_rel(rId)
    del xml_slides[index]


def load_template():
    buf = io.BytesIO()
    with zipfile.ZipFile(TMPL_PATH, 'r') as zin:
        with zipfile.ZipFile(buf, 'w', compression=zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename == '[Content_Types].xml':
                    data = data.replace(
                        b'presentationml.template.main+xml',
                        b'presentationml.presentation.main+xml',
                    )
                zout.writestr(item, data)
    buf.seek(0)
    return Presentation(buf)


# ── extract embedded images from ODP ────────────────────────────────────────

def extract_images():
    imgs = {}
    with zipfile.ZipFile(ODP_PATH, 'r') as z:
        for name in z.namelist():
            if name.startswith('media/'):
                imgs[name] = io.BytesIO(z.read(name))
    return imgs


# ── slide helpers ────────────────────────────────────────────────────────────

def _set_slide_number(slide, num):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 12:
            ph.text = str(num)


def _fill_text_frame(tf, bullets, font_size_pt=None):
    """Populate a text frame with a list of bullet items.
    Each item is either a str (level 0) or dict {'text': …, 'level': …}.
    """
    tf.word_wrap = True
    tf.clear()
    first = True
    for item in bullets:
        if isinstance(item, str):
            text, level = item, 0
        else:
            text, level = item['text'], item.get('level', 0)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.text  = text
        p.level = level
        if font_size_pt:
            for run in p.runs:
                run.font.size = Pt(font_size_pt)
        first = False


# Layout indices in the USC Engineering template:
#  0 Title Slide | 1 Title+Content | 2 Section Header
#  3 Two Content | 5 Title Only    | 7 Content w/ Caption

def add_title_slide(prs, title, authors, course_line, num):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = f"{authors}\n{course_line}"
    _set_slide_number(slide, num)
    return slide


def add_content_slide(prs, title, bullets, num, font_size_pt=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    _fill_text_frame(slide.placeholders[1].text_frame, bullets, font_size_pt)
    _set_slide_number(slide, num)
    return slide


def add_section_slide(prs, title, subtitle, num):
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    _set_slide_number(slide, num)
    return slide


def add_image_slide(prs, title, bullets, img_bytes_list, num):
    """Title-only layout with manual text box (left) + image(s) (right/tiled)."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    n_imgs = len(img_bytes_list)

    # ── text box ──
    if bullets:
        # For single-image slides: text on left third
        if n_imgs == 1:
            tx = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.9), Inches(5.0), Inches(5.2)
            )
        else:
            # For multi-image slides: text above images
            tx = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.7), Inches(12.3), Inches(1.1)
            )
        _fill_text_frame(tx.text_frame, bullets, font_size_pt=16)

    # ── image(s) ──
    if n_imgs == 1:
        img, _cap = img_bytes_list[0]
        img.seek(0)
        slide.shapes.add_picture(img, Inches(5.7), Inches(1.5), Inches(7.2), Inches(5.6))

    elif n_imgs == 2:
        for i, (img, _cap) in enumerate(img_bytes_list):
            img.seek(0)
            left = Inches(0.4 + i * 6.5)
            slide.shapes.add_picture(img, left, Inches(2.9), Inches(6.2), Inches(4.3))

    elif n_imgs == 3:
        for i, (img, _cap) in enumerate(img_bytes_list):
            img.seek(0)
            left = Inches(0.3 + i * 4.35)
            slide.shapes.add_picture(img, left, Inches(2.8), Inches(4.2), Inches(4.4))

    _set_slide_number(slide, num)
    return slide


# ── main ─────────────────────────────────────────────────────────────────────

def main():
    prs  = load_template()
    imgs = extract_images()

    # Remove sample slides that come baked into the template
    while len(prs.slides) > 0:
        _delete_slide(prs, 0)

    n    = 0

    def sn():
        nonlocal n
        n += 1
        return n

    # ── 1. Title ──────────────────────────────────────────────────────────────
    add_title_slide(
        prs,
        title       = "Sleep Stage Classification at the Edge",
        authors     = "Alex Anderson-McLeod  ·  Kuba Jerzmanowski  ·  Kasra Korminejad",
        course_line = "CSCE 714  ·  University of South Carolina",
        num         = sn(),
    )

    # ── 2. What Is Sleep Stage Classification? ────────────────────────────────
    add_content_slide(prs, "What Is Sleep Stage Classification?", [
        "Sleep cycles through four to five distinct, measurable physiological stages",
        "Stages per AASM: Wake, N1 (light), N2 (light), N3 (slow-wave/deep), REM",
        "Each stage exhibits characteristic EEG signatures — alpha waves, delta waves, sleep spindles, K-complexes",
        "Accurate staging reveals sleep architecture and its downstream effects on health and cognition",
    ], sn())

    # ── 3. Why Does It Matter? ────────────────────────────────────────────────
    add_content_slide(prs, "Why Does It Matter?", [
        "Sleep stage composition is a direct biomarker for downstream health outcomes",
        "Insufficient N3 (slow-wave sleep) \u2192 elevated Alzheimer\u2019s and metabolic disease risk",
        "REM deficiency \u2192 impaired memory consolidation and mood dysregulation",
        "Disrupted staging patterns \u2192 early indicator of sleep apnea, insomnia, narcolepsy",
        "Longitudinal staging enables clinical diagnosis and personalized treatment",
    ], sn())

    # ── 4. The Gold Standard: Clinical PSG ───────────────────────────────────
    add_content_slide(prs, "The Gold Standard: Clinical PSG", [
        "Polysomnography (PSG) \u2014 20+ channels: EEG, EOG, EMG, cardiorespiratory monitoring",
        "Multi-modal, full-head signal capture \u2192 high-fidelity staging across all AASM categories",
        "Automated classification accuracy: ~87\u201390% (5-class staging)",
        "Manual inter-rater Cohen\u2019s \u03ba \u2248 0.76\u20130.82 (AASM scoring guidelines)",
        "Requires a sleep clinic, trained technician, overnight stay \u2192 not viable for home use",
    ], sn())

    # ── 5. Smart Watches: Ubiquitous but Imprecise ────────────────────────────
    add_content_slide(prs, "Smart Watches: Ubiquitous but Imprecise", [
        "Now standard in consumer wearables \u2014 no additional hardware required",
        "Rely on PPG (photoplethysmography) and accelerometry as proxies for neural activity",
        "Cohen\u2019s \u03ba = 0.21\u20130.53 across six commercial devices [1]",
        "Particularly poor at distinguishing NREM substages (N1 vs. N2 vs. N3)",
        "Adequate for coarse sleep/wake detection; insufficient for clinically meaningful staging",
    ], sn())

    # ── 6. Our Approach: Wearable Single-Channel EEG ─────────────────────────
    add_content_slide(prs, "Our Approach: Wearable Single-Channel EEG", [
        "Custom lightweight headset \u2014 single EEG channel (two electrodes), designed for home use",
        "Captures direct neural signals \u2014 fundamentally more informative than wrist PPG",
        "Less intrusive than clinical PSG \u2014 practical for nightly, longitudinal monitoring",
        "Projected accuracy: ~80\u201386%, consistent with comparable single-channel EEG studies [2\u20134]",
        "Edge inference on Raspberry Pi \u2014 self-contained, no network or cloud dependency",
    ], sn())

    # ── 7. Why It Belongs on the Edge ────────────────────────────────────────
    add_content_slide(prs, "Why It Belongs on the Edge", [
        "Privacy \u2014 EEG contains sensitive biometric data; local processing avoids cloud transmission",
        "Connectivity independence \u2014 functions without a network connection",
        "Real-time latency \u2014 no round-trip delay to a remote server",
        "Accessibility \u2014 compact, self-contained form factor suitable for home deployment",
    ], sn())

    # ── 8. Available Sleep EEG Datasets ──────────────────────────────────────
    add_content_slide(prs, "Available Sleep EEG Datasets", [
        "SleepEDF \u2014 most widely used benchmark; annotated by AASM sleep stage",
        {"text": "Fpz-Cz probe locations (comfortable); two versions: SleepEDF-20 (20 subjects) and SleepEDF-78 (78 subjects)", "level": 1},
        "SHHS (Sleep Heart Health Study) \u2014 large-scale, but uses uncomfortable probe placement",
        "MASS (Montreal Archive of Sleep Studies) \u2014 requires institutional access review",
        "Selected: SleepEDF \u2014 openly accessible, widely benchmarked, compatible probe locations",
    ], sn())

    # ── 9. Section Divider: Prior Research ───────────────────────────────────
    add_section_slide(prs, "Prior Research", "Section 2", sn())

    # ── 10. Zhu et al. 2022 \u2014 intro ─────────────────────────────────────────────
    add_content_slide(prs, "Lightweight EEG Sleep Staging in Children at the Edge (Zhu et al., 2022)", [
        "Targets pediatric subjects \u2014 EEG staging patterns differ meaningfully from adults",
        "Emphasizes privacy benefits of edge-side EEG processing",
        "Five-class staging: Wake, N1, N2, N3, REM",
        "Datasets: SleepEDF (Fpz-Cz probes) and a custom pediatric set",
        "All experiments executed on a high-performance desktop \u2014 not deployed to edge hardware",
    ], sn())

    # ── 11. Zhu et al. 2022 \u2014 model + accuracy (image1) ────────────────────────
    imgs['media/image1.png'].seek(0)
    add_image_slide(prs,
        "Lightweight EEG Sleep Staging in Children at the Edge (Cont.)",
        [
            "Signal processing in the time domain (not frequency domain)",
            "Input segmented into 30-second epochs; normalized to zero mean and unit variance",
            "Accuracy: 83.06% (custom pediatric dataset), 86.41% (SleepEDF)",
            "Image credit: Zhu et al., World Wide Web, Dec. 2021.",
        ],
        [(imgs['media/image1.png'], "")],
        sn(),
    )

    # ── 12. Zhu et al. 2022 \u2014 architecture detail (image2) ────────────────────
    imgs['media/image2.png'].seek(0)
    add_image_slide(prs,
        "Lightweight EEG Sleep Staging in Children at the Edge (Cont.)",
        [
            "All activation functions: ReLU",
            "Image credit: Zhu et al., World Wide Web, Dec. 2021.",
        ],
        [(imgs['media/image2.png'], "")],
        sn(),
    )

    # ── 13. Atianashie Miracle et al. 2021 \u2014 intro ──────────────────────────────
    add_content_slide(prs, "CNN Sleep Disorder Classification on Raspberry Pi (Atianashie Miracle et al., 2021)", [
        "Binary classification: healthy vs. unhealthy (sleep apnea, insomnia)",
        "Deployed and tested on a Raspberry Pi \u2014 true edge inference",
        "Single EEG channel; signals processed in the time domain",
        "30-second epochs with noise filtering applied before classification",
        "Accuracy: 92% (sleep apnea), 89% (insomnia)",
    ], sn())

    # ── 14. CNN Pi \u2014 architecture (image3 + image4) ───────────────────────────
    imgs['media/image3.png'].seek(0)
    imgs['media/image4.png'].seek(0)
    add_image_slide(prs,
        "CNN Sleep Disorder Classification on Raspberry Pi (Cont.)",
        [
            "All activation functions: ReLU",
            "Image credit: Atianashie Miracle et al., J. Eng. Appl. Sci. Humanities, 2021.",
        ],
        [(imgs['media/image3.png'], ""), (imgs['media/image4.png'], "")],
        sn(),
    )

    # ── 15. Liu et al. 2023 \u2014 results (image5) ─────────────────────────────────
    imgs['media/image5.png'].seek(0)
    add_image_slide(prs,
        "MicroSleepNet: Mobile Real-Time Sleep Staging (Liu et al., 2023)",
        [
            "Deployed on Snapdragon-based Android; inference latency: 2.8 ms",
            "Evaluated on SleepEDF (Fpz-Cz) and SHHS (C4-A1)",
            "SleepEDF: no preprocessing; SHHS: class imbalance correction applied",
            "Standard 30-second epoch segmentation",
            "Accuracy: 83.3% (SHHS), 79.5% (SleepEDF-78), 82.8% (SleepEDF-20)",
            "Image credit: Liu et al., Frontiers in Neuroscience, Jul. 2023.",
        ],
        [(imgs['media/image5.png'], "")],
        sn(),
    )

    # ── 16. MicroSleepNet \u2014 architecture (image6 + image7 + image8) ──────────
    imgs['media/image6.png'].seek(0)
    imgs['media/image7.png'].seek(0)
    imgs['media/image8.png'].seek(0)
    add_image_slide(prs,
        "MicroSleepNet: Mobile Real-Time Sleep Staging (Cont.)",
        [
            "Channel shuffle applied between every two convolutional operations",
            "All activation functions: LeakyReLU",
            "Image credit: Liu et al., Frontiers in Neuroscience, Jul. 2023.",
        ],
        [
            (imgs['media/image6.png'], "DconvBlock"),
            (imgs['media/image7.png'], "GAP Block"),
            (imgs['media/image8.png'], ""),
        ],
        sn(),
    )

    # ── 17. References ────────────────────────────────────────────────────────
    add_content_slide(prs, "References", [
        "[1] A.-M. Schyvens et al., \u201cA performance validation of six commercial wrist-worn wearable sleep-tracking devices for sleep stage scoring compared to polysomnography,\u201d Sleep Advances, vol.\u00a06, no.\u00a02, Mar. 2025. doi: 10.1093/sleepadvances/zpaf021.",
        "[2] L. Zhu et al., \u201cA lightweight automatic sleep staging method for children using single-channel EEG based on edge artificial intelligence,\u201d World Wide Web, Dec. 2021. doi: 10.1007/s11280-021-00983-3.",
        "[3] A. Atianashie Miracle et al., \u201cA portable GUI based sleep disorder system classification based on convolution neural networks (CNN) in Raspberry Pi,\u201d J. Eng. Appl. Sci. Humanities, vol.\u00a06, pp.\u00a013\u201323, 2021.",
        "[4] G. Liu et al., \u201cMicroSleepNet: efficient deep learning model for mobile terminal real-time sleep staging,\u201d Frontiers in Neuroscience, vol.\u00a017, Jul. 2023. doi: 10.3389/fnins.2023.1218072.",
    ], sn(), font_size_pt=14)

    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}  ({n} slides)")


if __name__ == '__main__':
    main()
